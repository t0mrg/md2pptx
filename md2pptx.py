#!/usr/local/bin/python
# Script to convert Markdown to Powerpoint

import sys, getopt
from pptx import Presentation
from pptx.util import Inches


def processmd(infile,outfile):
    infilereader = open(infile,'r')
    indata = infilereader.readlines()
    infilereader.close()
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    for line in indata:
        line = line.strip()
        linetype = "na"
        if len(line) == 0:
            linetype = "ignore"
        else:
            if line[0] == "#":
                linetype = "h1"
                if line[0:2] == "##":
                    linetype = "h2"
                    if line[0:3] == "###":
                        linetype = "h3"
            if line.find("ncludegraphics")>0:
                linetype = "image"
        print line
        print linetype
        if linetype == "h1":
            slide = prs.slides.add_slide(title_slide_layout)
            shapes = slide.shapes
            title = slide.shapes.title
            title.text = line.replace("#","")
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
        if linetype == "h2":
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = line.replace("#","")
            title.text = line.replace("#","")
            tf = body_shape.text_frame
        if linetype == "h3":
            p = tf.add_paragraph()
            p.text = line.replace("#","")
            p.level = 1
        if linetype == "na":
            p = tf.add_paragraph()
            p.text = line.replace("-","")
            if line[0] == "\t":
                p.level = 3
            else:
                p.level = 2
        if linetype == "image":
            left = top = Inches(1)
            print line
            img_path = line.split("{")[1].replace("}","")
            pic = slide.shapes.add_picture(img_path, left, top)
    prs.save(outfile)

def main(argv):
    infile = ''
    outfile = ''
    try:
        opts, args = getopt.getopt(argv,"hi:o:",["ifile=","ofile="])
    except getopt.GetoptError:
        print 'Correct usage:\n    md2pptx.py -i <inputfile> -o <outputfile>'
        sys.exit(2)
    for opt, arg in opts:
        if opt == '-h':
            print 'Correct usage:\n    md2pptx.py -i <inputfile> -o <outputfile>'
            sys.exit()
        elif opt in ("-i", "--ifile"):
            infile = arg
        elif opt in ("-o", "--ofile"):
            outfile = arg
    processmd(infile,outfile)

if __name__ == "__main__":
   main(sys.argv[1:])
